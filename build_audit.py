#!/usr/bin/env python3
"""
Core logic for building RICOPA SEO audit deliverables (Excel + PowerPoint).
Importable module — no CLI dependency. Called by main.py (FastAPI).

v2: Auto-detects CSV/NDJSON files by column signatures.
    Falls back to filename hints for ambiguous cases.
    Supports mapping.json for training custom file→sheet assignments.

Usage:
    from build_audit import build_audit
    result = build_audit(data_dir="path/to/files",
                         xlsx="output.xlsx",
                         pptx="output.pptx",
                         client="Avafin",
                         domain="avafin.mx",
                         mapping={"mi_archivo.csv": "Errores 404"})
"""

from __future__ import annotations
import sys, os, json, re, shutil, csv
from pathlib import Path
from urllib.parse import urlparse
from datetime import date

import openpyxl
from openpyxl.styles import PatternFill, Font, Alignment
from openpyxl.utils import get_column_letter

from pptx import Presentation as PptxPresentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

BRAND_FONT = "Gabarito"
METADATA_FILL_A = PatternFill("solid", fgColor="FFF7E6")
METADATA_FILL_B = PatternFill("solid", fgColor="E6F4FF")


# ======================================================================
# Excel helpers
# ======================================================================

def _load_wb(path: str) -> openpyxl.Workbook:
    return openpyxl.load_workbook(path)

def _save_wb(wb, path: str) -> None:
    wb.save(path)

def _ensure_sheet(wb, name: str, headers: list[str]):
    if name in wb.sheetnames:
        ws = wb[name]
        ws.delete_rows(1, ws.max_row)
    else:
        ws = wb.create_sheet(title=name)
    ws.append(headers)
    for col_idx, _ in enumerate(headers, 1):
        ws.column_dimensions[get_column_letter(col_idx)].width = 22
    return ws

def _fill_sheet(wb, name: str, headers: list[str], rows: list[list]) -> int:
    ws = _ensure_sheet(wb, name, headers)
    n = 0
    for r in rows:
        ws.append(r)
        n += 1
    if n == 0:
        ws.append(["Sin hallazgos"] + [""] * (len(headers) - 1))
    return n


# ======================================================================
# Metadatos propuestos
# ======================================================================

def _fill_metadatos_propuestos(wb, rows: list[list]) -> None:
    name = "Metadatos propuestos"
    headers = ["URL", "Factor a corregir", "Texto actual", "Propuesta"]
    ws = _ensure_sheet(wb, name, headers)
    last_url = None
    toggle = False
    for r in rows:
        url = r[0]
        if url != last_url:
            toggle = not toggle
            last_url = url
        fill = METADATA_FILL_A if toggle else METADATA_FILL_B
        ws.append(r)
        row_idx = ws.max_row
        for col_idx in range(1, 5):
            cell = ws.cell(row=row_idx, column=col_idx)
            cell.fill = fill
            cell.alignment = Alignment(wrap_text=True, vertical="top")
    ws.column_dimensions["A"].width = 50
    ws.column_dimensions["B"].width = 28
    ws.column_dimensions["C"].width = 45
    ws.column_dimensions["D"].width = 60


# ======================================================================
# PowerPoint helpers
# ======================================================================

def _set_ppt_hallazgo(pptx_path: str, slide_index_1based: int, shape_substr: str,
                      count: int, excel_filename: str, sheet_name: str) -> int:
    prs = PptxPresentation(pptx_path)
    slide = prs.slides[slide_index_1based - 1]
    n_edited = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape_substr.lower() not in shape.name.lower():
            continue
        tf = shape.text_frame
        edited_here = 0
        for para in tf.paragraphs:
            for run in list(para.runs):
                t = run.text or ""
                if re.search(r"Hay \d+", t):
                    run.text = re.sub(r"Hay \d+", f"Hay {count}", t)
                    run.text = re.sub(r"\s*VER\s*$", "", run.text)
                    n_edited += 1
                    edited_here += 1
        if not edited_here:
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"Hay {count} hallazgos."
            n_edited += 1
            edited_here += 1
        last_para = tf.paragraphs[-1]
        ver_run = last_para.add_run()
        ver_run.text = "  VER"
        ver_run.font.name = BRAND_FONT
        ver_run.font.size = Pt(14)
        try:
            ver_run.font.color.rgb = RGBColor(0x05, 0x63, 0xC1)
        except Exception:
            pass
        try:
            ver_run.hyperlink.address = f"{excel_filename}#'{sheet_name}'!A1"
        except Exception:
            pass
    prs.save(pptx_path)
    return n_edited


# ======================================================================
# Column aliases: SF desktop CSVs use different names than MCP exports
# ======================================================================
COLUMN_ALIASES = {
    "desde": "Fuente",
    "hasta": "Destino",
    "código de respuesta": "Código de estado",
    "codigo de respuesta": "Código de estado",
    "tamaño": "Tamaño (bytes)",
    "transferido": "Código de estado",
    "rastreabilidad de enlaces": "Rastreabilidad",
    "tipo de ruta": "Tipo de ruta",
    "ruta del enlace": "Ruta del enlace",
    "origen del enlace": "Origen del enlace",
    "posición del enlace": "Posición del enlace",
    # NOT aliased: Texto ancla (different from Ancla in SF exports)
    "texto alt": "Texto ALT",
    "seguir": "Seguir",
    "destino": "Destino",
}

def _normalize_col(name: str) -> str:
    """Normalize column name using aliases."""
    return COLUMN_ALIASES.get(name.lower().strip(), name)


# ======================================================================
# Data processing
# ======================================================================

def _read_file(path: Path) -> list[dict]:
    """Auto-detect NDJSON or CSV, return list of dicts with normalized column names."""
    rows = []
    if not path.exists():
        return rows
    ext = path.suffix.lower()
    
    if ext == ".ndjson" or ext == ".jsonl":
        with open(path, encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    row = json.loads(line)
                    # Normalize keys
                    rows.append({_normalize_col(k): v for k, v in row.items()})
                except json.JSONDecodeError:
                    continue
    elif ext == ".csv":
        for encoding in ["utf-8-sig", "utf-8", "latin-1"]:
            try:
                with open(path, encoding=encoding, newline="") as f:
                    reader = csv.DictReader(f)
                    for row in reader:
                        clean = {}
                        for k, v in row.items():
                            norm_key = _normalize_col(k.strip() if k else k)
                            clean[norm_key] = v.strip() if isinstance(v, str) else v
                        rows.append(clean)
                if rows:
                    break
            except (UnicodeDecodeError, csv.Error):
                continue
    elif ext == ".json":
        with open(path, encoding="utf-8") as f:
            data = json.load(f)
            if isinstance(data, list):
                rows = [{_normalize_col(k): v for k, v in r.items()} for r in data if isinstance(r, dict)]
            elif isinstance(data, dict):
                rows = [{_normalize_col(k): v for k, v in data.items()}]
    return rows


def _get_columns(path: Path) -> frozenset:
    """Read the column headers from a file without loading all rows."""
    if not path.exists():
        return frozenset()
    ext = path.suffix.lower()
    
    if ext == ".ndjson":
        rows = _read_file(path)
        if rows and isinstance(rows[0], dict):
            return frozenset(rows[0].keys())
        return frozenset()
    
    if ext == ".csv":
        for encoding in ["utf-8-sig", "utf-8", "latin-1"]:
            try:
                with open(path, encoding=encoding, newline="") as f:
                    reader = csv.reader(f)
                    headers_raw = next(reader, [])
                    headers = frozenset(_normalize_col(h.strip()) for h in headers_raw if h.strip())
                    return headers
            except (UnicodeDecodeError, csv.Error):
                continue
        return frozenset()
    
    if ext == ".json":
        rows = _read_file(path)
        if rows and isinstance(rows[0], dict):
            return frozenset(rows[0].keys())
        return frozenset()
    
    return frozenset()


def _filter_by_domain(rows: list[dict], domain: str, key: str = "Dirección") -> list[dict]:
    if not domain:
        return rows
    result = []
    for r in rows:
        url = r.get(key, "")
        if not url or not isinstance(url, str):
            continue
        host = urlparse(url).hostname or ""
        if host.endswith(domain) or host == domain:
            result.append(r)
    return result


def _dedup_by_key(rows: list[dict], key: str = "Dirección") -> list[dict]:
    seen = set()
    result = []
    for r in rows:
        val = r.get(key, "")
        if val in seen:
            continue
        seen.add(val)
        result.append(r)
    return result


# ======================================================================
# Data sampling: detect sheet by reading first rows' values
# ======================================================================

def _sample_detect(fpath: Path) -> str | None:
    """
    Read the first 10 rows and, based on the actual data values,
    determine which sheet this file belongs to.
    Only called when column signature is ambiguous or unknown.
    Returns sheet name or None.
    """
    rows = _read_file(fpath)
    if not rows:
        return None
    cols = _get_columns(fpath)
    sample = rows[:10]
    
    def _int_val(r, field):
        try:
            return int(float(r.get(field, 0) or 0))
        except (ValueError, TypeError):
            return 0
    
    # --- Title files (Dirección, Título 1, Longitud, Ancho) ---
    if frozenset(["Dirección", "Título 1"]) <= cols:
        longs = [_int_val(r, "Longitud del título 1") for r in sample if _int_val(r, "Longitud del título 1") > 0]
        reps = [_int_val(r, "Repeticiones") for r in sample]
        titulos = [r.get("Título 1","") for r in sample if r.get("Título 1","").strip()]
        
        if not titulos and all(v == 0 for v in longs):
            return "Metatítulo — Falta"
        if longs and all(v > 60 for v in longs):
            return "Metatítulo — Largo"
        if longs and all(v < 30 for v in longs):
            return "Metatítulo — Corto"
        if any(v > 0 for v in reps) and all(v > 1 for v in reps if v > 0):
            return "Metatítulo — Duplicado"
        # Mixed lengths → likely Todo export, handled by derivation
        if longs and any(v > 60 for v in longs) and any(v < 30 for v in longs):
            return None  # Let derivation handle it
        if longs:
            return "Metatítulo — Largo"  # Default: most common
    
    # --- Meta files (Dirección, Meta description 1, Longitud, Ancho) ---
    if frozenset(["Dirección", "Meta description 1"]) <= cols:
        longs = [_int_val(r, "Longitud de la meta description 1") for r in sample if _int_val(r, "Longitud de la meta description 1") > 0]
        reps = [_int_val(r, "Repeticiones") for r in sample]
        metas = [r.get("Meta description 1","") for r in sample if r.get("Meta description 1","").strip()]
        
        if not metas and all(v == 0 for v in longs):
            return "Metadescripción — Falta"
        if longs and all(v > 155 for v in longs):
            return "Metadescripción — Larga"
        if longs and all(v < 70 for v in longs):
            return "Metadescripción — Corta"
        if any(v > 0 for v in reps) and all(v > 1 for v in reps if v > 0):
            return "Metadescripción — Duplicada"
        return None  # Todo export
    
    # --- H1 files (Dirección, H1-1, Longitud) ---
    if frozenset(["Dirección", "H1-1"]) <= cols:
        longs = [_int_val(r, "Longitud de H1-1") for r in sample if _int_val(r, "Longitud de H1-1") > 0]
        reps = [_int_val(r, "Repeticiones") for r in sample]
        h1s = [r.get("H1-1","") for r in sample if r.get("H1-1","").strip()]
        
        if not h1s and all(v == 0 for v in longs):
            return "H1 — Falta"
        if longs and all(v > 70 for v in longs):
            return "H1 — Largo (+70)"
        if any(v > 0 for v in reps) and all(v > 1 for v in reps if v > 0):
            return "H1 — Duplicado"
        return None  # Todo export
    
    # --- Images: +100kB ---
    if frozenset(["Dirección", "Tamaño (Bytes)"]) <= cols:
        sizes = [_int_val(r, "Tamaño (Bytes)") for r in sample if _int_val(r, "Tamaño (Bytes)") > 0]
        if sizes and all(s > 100000 for s in sizes):
            return "Imágenes +100kb"
        return None
    
    # --- Images: size attr missing ---
    if frozenset(["Dirección", "Dimensiones"]) <= cols:
        dims = [r.get("Dimensiones") for r in sample if r.get("Dimensiones")]
        if not dims:
            return "Imágenes sin atributo tamaño"
        return None
    
    # --- Images: ALT missing (both Falta texto ALT and Falta atributo ALT) ---
    if frozenset(["Dirección", "Texto ALT"]) <= cols:
        alts = [r.get("Texto ALT", "") for r in sample if r.get("Texto ALT", "").strip()]
        if not alts:
            return "Imágenes sin ALT text"
    
    # --- Images: ALT from SEO element export (Dirección, Tipo, Enlaces de entrada) ---
    if frozenset(["Dirección", "Enlaces de entrada de imágenes"]) <= cols:
        return "Imágenes sin ALT text"
        return None
    
    # --- Canonicals ---
    if frozenset(["Dirección", "Elemento de enlace canónico 1"]) <= cols:
        can = [r.get("Elemento de enlace canónico 1") for r in sample if r.get("Elemento de enlace canónico 1", "").strip()]
        if not can:
            return "Canónica — Falta"
        return "Canónica — Múltiple"
    
    # --- Bulk exports: Fuente-based files ---
    if frozenset(["Fuente", "Destino"]) <= cols:
        # Check if the data suggests images (image URLs in Destino)
        destinos = [r.get("Destino", "") for r in sample if r.get("Destino")]
        if destinos:
            first = destinos[0].lower()
            if any(ext in first for ext in [".jpg", ".png", ".webp", ".gif", ".svg"]):
                # Image-related bulk export
                tipo = [r.get("Tipo", "") for r in sample]
                if tipo and "Imagen" in str(tipo):
                    if "Tamaño (bytes)" in cols:
                        return "Detalle Imágenes +100kb"
                    return "Detalle Imágenes sin ALT"
            return None  # Could be canonical, noindex, etc.
    
    # --- PS reports ---
    if frozenset(["Página fuente"]) <= cols:
        # Report files — check URL patterns and metrics
        urls = [r.get("URL", "") for r in sample if r.get("URL")]
        if urls:
            first = urls[0].lower()
            if any(ext in first for ext in [".js", ".css"]):
                if "Posible ahorro (bytes)" in cols:
                    return "PS Minificar JS"
                return "PS Minificar CSS"
            if any(ext in first for ext in [".woff", ".woff2", ".ttf", ".otf"]):
                return "PS Visualización fuentes"
            if any(ext in first for ext in [".jpg", ".png", ".webp", ".gif", ".svg"]):
                if "Motivo" in cols:
                    return "PS Mejorar entrega imágenes"
            return None
    
    # --- Images Todo (mixed data check) ---
    if frozenset(["Dirección", "Dimensiones", "Tipo de contenido", "Enlaces de entrada de imágenes"]).issubset(cols):
        sizes = [_int_val(r, "Tamaño (Bytes)") for r in sample if _int_val(r, "Tamaño (Bytes)") > 0]
        if sizes:
            over = [s for s in sizes if s > 100000]
            under = [s for s in sizes if s <= 100000]
            if over and under:
                return "__derive_images"
        # Without size, check if mixed ALT/Dimensions
        alts_good = [r for r in sample if (r.get("Texto ALT") or "").strip()]
        alts_bad = [r for r in sample if not (r.get("Texto ALT") or "").strip()]
        dims_good = [r for r in sample if r.get("Dimensiones")]
        dims_bad = [r for r in sample if not r.get("Dimensiones")]
        has_mix = (alts_good and alts_bad) or (dims_good and dims_bad)
        if has_mix:
            return "__derive_images"
        return "Imágenes +100kb"  # All large → filter export
    
    return None


# ======================================================================
# Sub-sheet derivation from "Todo" exports
# ======================================================================

# Element type detection by column fingerprint
ELEMENT_FINGERPRINTS = {
    "title": frozenset(["Dirección", "Título 1", "Longitud del título 1", "Ancho de píxeles del título 1"]),
    "meta": frozenset(["Dirección", "Meta description 1", "Longitud de la meta description 1", "Ancho de píxeles de la meta description 1"]),
    "h1": frozenset(["Dirección", "H1-1", "Longitud de H1-1"]),
    "images": frozenset(["Dirección", "Dimensiones", "Tipo de contenido", "Enlaces de entrada de imágenes"]),
}

def _detect_element_type(columns: frozenset) -> str | None:
    """Detect if columns belong to a known 'Todo' export element type."""
    for elem_type, fingerprint in ELEMENT_FINGERPRINTS.items():
        if fingerprint.issubset(columns):
            return elem_type
    return None

def _derive_subsheets(rows: list[dict], element_type: str) -> dict:
    """
    From a 'Todo' export of an element, derive all sub-sheets.
    Returns {sheet_name: (headers, rows)}.
    """
    results = {}
    
    if element_type == "title":
        # Metatítulo — Falta
        falta = [r for r in rows if not (r.get("Título 1") or "").strip()]
        if falta:
            results["Metatítulo — Falta"] = (
                ["Dirección", "Indexabilidad"],
                [[r.get("Dirección",""), r.get("Indexabilidad","")] for r in falta]
            )
        # Metatítulo — Duplicado
        dup = [r for r in rows if int(r.get("Repeticiones", 0) or 0) > 1]
        if dup:
            results["Metatítulo — Duplicado"] = (
                ["Dirección", "Título", "Repeticiones"],
                [[r.get("Dirección",""), r.get("Título 1",""), r.get("Repeticiones","")] for r in dup]
            )
        # Metatítulo — Largo (>60 chars or >561 px)
        largo = [r for r in rows if (int(r.get("Longitud del título 1", 0) or 0) > 60 or 
                                      int(r.get("Ancho de píxeles del título 1", 0) or 0) > 561)]
        if largo:
            results["Metatítulo — Largo"] = (
                ["Dirección", "Título", "Caracteres", "Píxeles"],
                [[r.get("Dirección",""), r.get("Título 1",""), 
                  r.get("Longitud del título 1",""), r.get("Ancho de píxeles del título 1","")] for r in largo]
            )
        # Metatítulo — Corto (<30 chars or <200 px)
        corto = [r for r in rows if (int(r.get("Longitud del título 1", 999) or 999) < 30 or 
                                      int(r.get("Ancho de píxeles del título 1", 999) or 999) < 200)]
        if corto:
            results["Metatítulo — Corto"] = (
                ["Dirección", "Título", "Caracteres", "Píxeles"],
                [[r.get("Dirección",""), r.get("Título 1",""),
                  r.get("Longitud del título 1",""), r.get("Ancho de píxeles del título 1","")] for r in corto]
            )
    
    elif element_type == "meta":
        # Metadescripción — Falta
        falta = [r for r in rows if not (r.get("Meta description 1") or "").strip()]
        if falta:
            results["Metadescripción — Falta"] = (
                ["Dirección", "Indexabilidad"],
                [[r.get("Dirección",""), r.get("Indexabilidad","")] for r in falta]
            )
        # Metadescripción — Duplicada
        dup = [r for r in rows if int(r.get("Repeticiones", 0) or 0) > 1]
        if dup:
            results["Metadescripción — Duplicada"] = (
                ["Dirección", "Metadescripción", "Repeticiones"],
                [[r.get("Dirección",""), r.get("Meta description 1",""), r.get("Repeticiones","")] for r in dup]
            )
        # Metadescripción — Larga (>155 chars or >985 px)
        largo = [r for r in rows if (int(r.get("Longitud de la meta description 1", 0) or 0) > 155 or
                                      int(r.get("Ancho de píxeles de la meta description 1", 0) or 0) > 985)]
        if largo:
            results["Metadescripción — Larga"] = (
                ["Dirección", "Metadescripción", "Caracteres", "Píxeles"],
                [[r.get("Dirección",""), r.get("Meta description 1",""),
                  r.get("Longitud de la meta description 1",""), r.get("Ancho de píxeles de la meta description 1","")] for r in largo]
            )
        # Metadescripción — Corta (<70 chars or <400 px)
        corto = [r for r in rows if (int(r.get("Longitud de la meta description 1", 999) or 999) < 70 or
                                      int(r.get("Ancho de píxeles de la meta description 1", 999) or 999) < 400)]
        if corto:
            results["Metadescripción — Corta"] = (
                ["Dirección", "Metadescripción", "Caracteres", "Píxeles"],
                [[r.get("Dirección",""), r.get("Meta description 1",""),
                  r.get("Longitud de la meta description 1",""), r.get("Ancho de píxeles de la meta description 1","")] for r in corto]
            )
    
    elif element_type == "h1":
        # H1 — Falta
        falta = [r for r in rows if not (r.get("H1-1") or "").strip()]
        if falta:
            results["H1 — Falta"] = (
                ["Dirección", "Indexabilidad"],
                [[r.get("Dirección",""), r.get("Indexabilidad","")] for r in falta]
            )
        # H1 — Duplicado
        dup = [r for r in rows if int(r.get("Repeticiones", 0) or 0) > 1]
        if dup:
            results["H1 — Duplicado"] = (
                ["Dirección", "H1", "Repeticiones"],
                [[r.get("Dirección",""), r.get("H1-1",""), r.get("Repeticiones","")] for r in dup]
            )
        # H1 — Largo (>70 chars)
        largo = [r for r in rows if int(r.get("Longitud de H1-1", 0) or 0) > 70]
        if largo:
            results["H1 — Largo (+70)"] = (
                ["Dirección", "H1", "Caracteres"],
                [[r.get("Dirección",""), r.get("H1-1",""), r.get("Longitud de H1-1","")] for r in largo]
            )
    
    elif element_type == "images":
        # Imágenes — +100kb (skip if Tamaño column not present)
        col_set = set(rows[0].keys()) if rows else set()
        if "Tamaño (Bytes)" in col_set:
            over100 = [r for r in rows if int(r.get("Tamaño (Bytes)", 0) or 0) > 100000]
            if over100:
                results["Imágenes +100kb"] = (
                    ["Dirección", "Tipo", "Tamaño (kB)", "Dimensiones", "# enlaces entrada"],
                    [[r.get("Dirección",""), r.get("Tipo de contenido",""), r.get("Tamaño (Bytes)",""),
                      r.get("Dimensiones",""), r.get("Enlaces de entrada de imágenes","")] for r in over100]
                )
        # Imágenes — sin ALT (both Falta texto ALT and Falta atributo ALT go here)
        sin_alt = [r for r in rows if not (r.get("Texto ALT") or "").strip()]
        if sin_alt:
            results["Imágenes sin ALT text"] = (
                ["Dirección", "Tipo", "# páginas donde aparece"],
                [[r.get("Dirección",""), r.get("Tipo de contenido",""),
                  r.get("Enlaces de entrada de imágenes","")] for r in sin_alt]
            )
        # Imágenes — sin atributo tamaño
        sin_size = [r for r in rows if not r.get("Dimensiones") or str(r.get("Dimensiones","")).lower() in ("", "none")]
        if sin_size:
            results["Imágenes sin atributo tamaño"] = (
                ["Dirección", "Tipo", "Dimensiones"],
                [[r.get("Dirección",""), r.get("Tipo de contenido",""), r.get("Dimensiones","")] for r in sin_size]
            )
    
    return results


# ======================================================================
# Column signatures → sheet name detection
# ======================================================================
# Each entry: signature (frozenset of columns) → (sheet_name, headers, col_map)
# col_map maps source columns → output columns (None = skip)

SIGNATURES = {}

def _register(columns, sheet_name, headers, col_map):
    key = frozenset(columns)
    # Multiple sheets can have the same signature (e.g., all "Falta" exports share Dirección,Indexabilidad)
    entry = {"sheet": sheet_name, "headers": headers, "col_map": col_map}
    if key not in SIGNATURES:
        SIGNATURES[key] = []
    SIGNATURES[key].append(entry)

# --- Signatures from SEO element exports ---
_register(["Dirección", "Respuesta", "Tipo de contenido", "URL de redirección", "Tipo de redirección"],
          "Errores 404", 
          ["URL", "Código", "Tipo Contenido", "URL de redirección", "Tipo redirección"],
          ["Dirección", "Respuesta", "Tipo de contenido", "URL de redirección", "Tipo de redirección"])

_register(["Dirección", "Respuesta", "URL de redirección", "Tipo de redirección", "Enlaces internos"],
          "Redirecciones 3xx",
          ["URL", "Código", "URL destino", "Tipo", "Enlaces internos"],
          ["Dirección", "Respuesta", "URL de redirección", "Tipo de redirección", "Enlaces internos"])

_register(["Dirección", "Indexabilidad"],  # ambiguous: title_falta, meta_falta, h1_falta, can_falta
          None, None, None)  # resolved by filename hints

# Individual "Falta" sheets all share the same columns, registered for hint-based detection
_register(["Dirección", "Indexabilidad"], "Metatítulo — Falta", ["URL", "Indexabilidad"], ["Dirección", "Indexabilidad"])
_register(["Dirección", "Indexabilidad"], "Metadescripción — Falta", ["URL", "Indexabilidad"], ["Dirección", "Indexabilidad"])
_register(["Dirección", "Indexabilidad"], "H1 — Falta", ["URL", "Indexabilidad"], ["Dirección", "Indexabilidad"])
_register(["Dirección", "Indexabilidad"], "Canónica — Falta", ["URL", "Indexabilidad"], ["Dirección", "Indexabilidad"])
_register(["Dirección", "Indexabilidad"], "Poco contenido", ["URL", "Indexabilidad"], ["Dirección", "Indexabilidad"])
_register(["Dirección", "Indexabilidad"], "URLs no ASCII", ["URL", "Indexabilidad"], ["Dirección", "Indexabilidad"])

_register(["Dirección"],  # data-only export, no indexability info
          None, None, None)  # resolved by filename hints

# Individual Dirección-only entries for hint matching
_register(["Dirección"], "Canónica — Múltiple", ["URL", "Canónicas"], ["Dirección", None])
_register(["Dirección"], "Canónica — Errores", ["URL", "Tipo problema", "Canónica", "Estado"], ["Dirección", None])
_register(["Dirección"], "PS Minificar JS", ["URL", "Ahorro (ms)"], ["Dirección", None])
_register(["Dirección"], "PS Minificar CSS", ["URL", "Ahorro (ms)"], ["Dirección", None])
_register(["Dirección"], "PS Visualización fuentes", ["URL", "Ahorro (ms)"], ["Dirección", None])
_register(["Dirección"], "PS Mejorar entrega imágenes", ["URL", "URL imagen", "Ahorro"], ["Dirección", None])
_register(["Dirección"], "PS Solicitudes LCP", ["URL", "Elemento", "Tipo LCP"], ["Dirección", None])
_register(["Dirección"], "PageSpeed — CLS", ["URL", "Elemento", "Puntuación CLS"], ["Dirección", None])

_register(["Dirección", "Título 1", "Longitud del título 1", "Ancho de píxeles del título 1"],
          "Metatítulo — Largo",
          ["URL", "Título", "Caracteres", "Píxeles"],
          ["Dirección", "Título 1", "Longitud del título 1", "Ancho de píxeles del título 1"])

_register(["Dirección", "Título 1", "Longitud del título 1", "Ancho de píxeles del título 1"],
          "Metatítulo — Corto",
          ["URL", "Título", "Caracteres", "Píxeles"],
          ["Dirección", "Título 1", "Longitud del título 1", "Ancho de píxeles del título 1"])

_register(["Dirección", "Título 1", "Repeticiones"],
          "Metatítulo — Duplicado",
          ["URL", "Título", "Repeticiones"],
          ["Dirección", "Título 1", "Repeticiones"])

_register(["Dirección", "Meta description 1", "Longitud de la meta description 1", "Ancho de píxeles de la meta description 1"],
          "Metadescripción — Larga",
          ["URL", "Metadescripción", "Caracteres", "Píxeles"],
          ["Dirección", "Meta description 1", "Longitud de la meta description 1", "Ancho de píxeles de la meta description 1"])

_register(["Dirección", "Meta description 1", "Repeticiones"],
          "Metadescripción — Duplicada",
          ["URL", "Metadescripción", "Repeticiones"],
          ["Dirección", "Meta description 1", "Repeticiones"])

_register(["Dirección", "H1-1", "Repeticiones"],
          "H1 — Duplicado",
          ["URL", "H1", "Repeticiones"],
          ["Dirección", "H1-1", "Repeticiones"])

_register(["Dirección", "H1-1", "Longitud de H1-1"],
          "H1 — Largo (+70)",
          ["URL", "H1", "Caracteres"],
          ["Dirección", "H1-1", "Longitud de H1-1"])

_register(["Dirección", "H1-1"],
          None, None, None)  # ambiguous: h1_multi or h1_largo page export, use filename hint

_register(["Dirección", "Tipo de contenido", "Tamaño (Bytes)", "Dimensiones", "Enlaces de entrada de imágenes"],
          "Imágenes +100kb",
          ["URL imagen", "Tipo", "Tamaño (kB)", "Dimensiones", "# enlaces entrada"],
          ["Dirección", "Tipo de contenido", "Tamaño (Bytes)", "Dimensiones", "Enlaces de entrada de imágenes"])

_register(["Dirección", "Tipo de contenido", "Enlaces de entrada de imágenes"],
          "Imágenes sin ALT text",
          ["URL imagen", "Tipo", "# páginas donde aparece"],
          ["Dirección", "Tipo de contenido", "Enlaces de entrada de imágenes"])

_register(["Dirección", "Tipo de contenido", "Dimensiones"],
          "Imágenes sin atributo tamaño",
          ["URL imagen", "Tipo", "Dimensiones"],
          ["Dirección", "Tipo de contenido", "Dimensiones"])

_register(["Dirección", "Elemento de enlace canónico 1"],
          "Canónica — Múltiple",
          ["URL", "Canónicas"],
          ["Dirección", "Elemento de enlace canónico 1"])

_register(["Dirección", "Indexabilidad", "Estado de indexabilidad"],
          "Canónica — Errores",
          ["URL", "Tipo problema", "Canónica", "Estado"],
          ["Dirección", None, "Elemento de enlace canónico 1", "Estado de indexabilidad"])

# --- Signatures from bulk exports ---
_register(["Tipo", "Fuente", "Destino", "Tamaño (bytes)", "Ancla", "Código de estado"],
          None, None, None)  # generic bulk export, resolved by specific subsets

_register(["Fuente", "Destino", "Tamaño (bytes)"],
          "Detalle Imágenes +100kb",
          ["URL imagen", "Tamaño (kB)", "URL donde aparece"],
          ["Destino", "Tamaño (bytes)", "Fuente"])

_register(["Fuente", "Destino", "Texto ALT", "Ancla"],
          "Detalle Imágenes sin ALT",
          ["URL imagen", "URL página", "Texto ancla"],
          ["Destino", "Fuente", "Ancla"])

_register(["Fuente", "Destino"],
          "Detalle Imágenes sin size attr",
          ["URL imagen", "URL página"],
          ["Destino", "Fuente"])

_register(["Fuente", "Destino", "Texto ALT", "Longitud", "Tipo de ruta", "Posición del enlace", "Origen del enlace"],
          None, None, None)  # ambiguous: noindex/nofollow, resolve by dir name

_register(["Fuente", "Destino", "Ancla"],
          None, None, None)  # ambiguous
    
# Detalle Errores 404 from SF desktop bulk export (Desde/Hasta → Fuente/Destino)
_register(["Fuente", "Destino", "Texto ancla", "Código de estado"],
          "Detalle Errores 404",
          ["URL 404", "URL Origen", "Texto ancla"],
          ["Destino", "Fuente", "Texto ancla"])

# Detalle 301 redirects from SF desktop bulk export
_register(["Fuente", "Destino", "Follow", "Código de estado", "Origen del enlace"],
          "Redirecciones 3xx",
          ["URL", "Código", "URL destino"],
          ["Fuente", "Código de estado", "Destino"])

# --- Signatures from PageSpeed SEO element exports (Dirección-based, aggregate) ---
_register(["Dirección", "Ahorro al minimizar JavaScript (ms)", "Ahorro al minimizar JavaScript (Bytes)"],
          "PS Minificar JS",
          ["URL", "Ahorro (ms)", "Ahorro estimado"],
          ["Dirección", "Ahorro al minimizar JavaScript (ms)", "Ahorro al minimizar JavaScript (Bytes)"])

_register(["Dirección", "Ahorro al minimizar CSS (ms)"],
          "PS Minificar CSS",
          ["URL", "Ahorro (ms)"],
          ["Dirección", "Ahorro al minimizar CSS (ms)"])

_register(["Dirección", "Ahorro en la visualización de fuentes (ms)"],
          "PS Visualización fuentes",
          ["URL", "Ahorro (ms)"],
          ["Dirección", "Ahorro en la visualización de fuentes (ms)"])

# --- Signatures from PageSpeed reports (file-level detail) ---
_register(["Página fuente", "URL", "Tamaño (bytes)", "Posible ahorro (bytes)"],
          "PS Minificar JS",
          ["URL Página", "Archivo JS a minificar", "Tamaño (bytes)", "Ahorro estimado (bytes)"],
          ["Página fuente", "URL", "Tamaño (bytes)", "Posible ahorro (bytes)"])

# PageSpeed reports from SF desktop (Spanish column names)
_register(["Página fuente", "URL", "Tamaño (bytes)", "Posible ahorro (bytes)", "Motivo"],
          "PS Mejorar entrega imágenes",
          ["URL Página", "Imagen a optimizar", "Tamaño (bytes)", "Ahorro (bytes)", "Motivo"],
          ["Página fuente", "URL", "Tamaño (bytes)", "Posible ahorro (bytes)", "Motivo"])

_register(["Página fuente", "URL", "Posible ahorro (bytes)", "Motivo"],
          "PS Mejorar entrega imágenes",
          ["URL Página", "Imagen a optimizar", "Tamaño (bytes)", "Ahorro (bytes)", "Motivo"],
          ["Página fuente", "URL", "Tamaño (bytes)", "Posible ahorro (bytes)", "Motivo"])

_register(["Página fuente", "URL", "Posible ahorro (ms)"],
          "PS Visualización fuentes",
          ["URL Página", "Archivo de fuente", "Ahorro (ms)"],
          ["Página fuente", "URL", "Posible ahorro (ms)"])

_register(["Página fuente", "Selector", "Carga diferida no aplicada", 'Se ha de aplicar "fetchpriority=high"', "La solicitud es detectable en el documento inicial"],
          "PS Solicitudes LCP",
          ["URL Página", "Selector LCP", "Carga diferida", "fetchpriority=high", "Detectable inicial"],
          ["Página fuente", "Selector", "Carga diferida no aplicada", 'Se ha de aplicar "fetchpriority=high"', "La solicitud es detectable en el documento inicial"])

_register(["Página fuente", "Etiqueta", "Contribución de CLS", "Recortes"],
          "PageSpeed — CLS",
          ["URL Página", "Elemento causante", "Contribución CLS", "Fragmento HTML"],
          ["Página fuente", "Etiqueta", "Contribución de CLS", "Recortes"])

# --- Detalle 404 ---
_register(["URL 404", "URL Origen", "Texto ancla"],
          "Detalle Errores 404",
          ["URL 404", "URL Origen", "Texto ancla"],
          ["URL 404", "URL Origen", "Texto ancla"])


# Filename hints for ambiguous signatures (Dirección + Indexabilidad, Fuente + Destino + Ancla, etc.)
FILENAME_HINTS = {
    "title_falta": "Metatítulo — Falta",
    "title_dup": "Metatítulo — Duplicado",
    "title_multiple": "Metatítulo — Múltiple",
    "title_largo": "Metatítulo — Largo",
    "title_corto": "Metatítulo — Corto",
    "title_h1": "Metatítulo = H1",
    "meta_falta": "Metadescripción — Falta",
    "meta_dup": "Metadescripción — Duplicada",
    "meta_larga": "Metadescripción — Larga",
    "meta_corta": "Metadescripción — Corta",
    "meta_multiple": "Metadescripción — Múltiple",
    "h1_falta": "H1 — Falta",
    "h1_dup": "H1 — Duplicado",
    "h1_multiple": "H1 — Múltiple",
    "h1_largo": "H1 — Largo (+70)",
    "can_falta": "Canónica — Falta",
    "can_multiple": "Canónica — Múltiple",
    "can_errores": "Canónica — Errores",
    "canónica_múltiple": "Canónica — Múltiple",
    "canónica_errores": "Canónica — Errores",
    "poco_contenido": "Poco contenido",
    "urls_no_ascii": "URLs no ASCII",
    "urls_http": "URLs HTTP (no HTTPS)",
    "img_over_100kb": "Imágenes +100kb",
    "img_over_100kb_detalle": "Detalle Imágenes +100kb",
    "img_sin_alt": "Imágenes sin ALT text",
    "img_sin_alt_detalle": "Detalle Imágenes sin ALT",
    "img_sin_size": "Imágenes sin atributo tamaño",
    "img_sin_size_detalle": "Detalle Imágenes sin size attr",
    "img_dim_incorrectas": "Imágenes dim incorrectas",
    "img_dim_incorrectas_detalle": "Detalle Img dim incorrectas",
    "noindex": "Directivas — Noindex",
    "nofollow": "Directivas — Nofollow",
    "errores_404": "Errores 404",
    "errors_4xx": "Errores 404",
    "detalle_errores_404": "Detalle Errores 404",
    "redirects_3xx": "Redirecciones 3xx",
    "redireccion": "Redirecciones 3xx",
    "bucles_cadena": "Bucles y cadenas redirección",
    "ps_minificar_js": "PS Minificar JS",
    "ps_minificar_css": "PS Minificar CSS",
    "ps_fuentes": "PS Visualización fuentes",
    "ps_visualizacion": "PS Visualización fuentes",
    "ps_visualización": "PS Visualización fuentes",
    "ps_mejorar_img": "PS Mejorar entrega imágenes",
    "ps_mejorar_entrega": "PS Mejorar entrega imágenes",
    "ps_lcp": "PS Solicitudes LCP",
    "ps_solicitudes_lcp": "PS Solicitudes LCP",
    "ps_cls": "PageSpeed — CLS",
    "pagespeed_cls": "PageSpeed — CLS",
    "ps_report_fuentes": "PS Visualización fuentes",
    "ps_report_minificar": "PS Minificar JS",
    "ps_report_mejorar": "PS Mejorar entrega imágenes",
    "ps_report_lcp": "PS Solicitudes LCP",
    "ps_report_cls": "PageSpeed — CLS",
}


def _hint_from_filename(fname: str) -> str | None:
    """Try to guess sheet name from filename substring hints."""
    fname_lower = fname.lower().replace(".csv", "").replace(".ndjson", "").replace(" ", "_").replace("-", "_")
    for hint, sheet in FILENAME_HINTS.items():
        if hint in fname_lower:
            return sheet
    return None


def _detect_assignments(data_dir: str, mapping: dict = None) -> tuple[dict[str, tuple[list[str], list[str], str]], list[dict]]:
    """
    Scan data_dir for CSV/NDJSON files, detect their sheet assignment.
    
    Returns: (assignments, unmatched)
        assignments: sheet_name → (headers, col_map, file_path)
        unmatched: [{file, columns, hint}] for files that couldn't be auto-assigned
    """
    data_dir = Path(data_dir)
    mapping = mapping or {}
    assignments = {}
    unmatched = []
    processed_files = set()  # filenames already assigned
    
    # Build reverse mapping: filename → sheet_name from mapping
    reverse_map = {}
    for fname, sheet in (mapping or {}).items():
        reverse_map[fname.strip().lower()] = sheet
    
    # Find all data files
    data_files = []
    for ext in [".csv", ".ndjson", ".json"]:
        data_files.extend(data_dir.glob(f"*{ext}"))
    # Flatten: if files are in subdirs, move consideration to root
    for root, dirs, fnames in os.walk(data_dir):
        for fname in fnames:
            fpath = Path(root) / fname
            if fpath not in data_files and fpath.suffix.lower() in [".csv", ".ndjson", ".json"]:
                data_files.append(fpath)
    
    for fpath in sorted(data_files):
        fname = fpath.name
        if fname.startswith(".") or fname.startswith("__"):
            continue
        
        cols = _get_columns(fpath)
        if not cols:
            continue
        
        assignment = None
        
        # 1. Check mapping.json
        for key, sheet in reverse_map.items():
            if key in fname.lower():
                # Find the signature for this sheet
                for sig, entries in SIGNATURES.items():
                    for e in entries:
                        if e["sheet"] == sheet:
                            assignment = (sheet, e["headers"], e["col_map"], str(fpath))
                            break
        
        # 2. Check column signatures (subset match: file columns must CONTAIN the signature columns)
        if not assignment:
            entries = []
            for sig_key, sig_entries in SIGNATURES.items():
                if sig_key.issubset(cols):
                    entries.extend([(len(sig_key), e) for e in sig_entries if e["sheet"] is not None])
            # Pick the most specific match (longest signature = most columns matched)
            if entries:
                entries.sort(key=lambda x: -x[0])  # Descending by specificity
                best_len = entries[0][0]
                # Get all entries at the same specificity level
                best_matches = [e for l, e in entries if l == best_len]
                if len(best_matches) == 1:
                    e = best_matches[0]
                    assignment = (e["sheet"], e["headers"], e["col_map"], str(fpath))
                elif len(best_matches) > 1:
                    # Same specificity but multiple matches — try filename hint
                    hint = _hint_from_filename(fname)
                    if hint:
                        for e in best_matches:
                            if e["sheet"] == hint:
                                assignment = (e["sheet"], e["headers"], e["col_map"], str(fpath))
                                break
        
        # 3. Data sampling: read actual values to determine sheet
        if not assignment:
            detected = _sample_detect(fpath)
            if detected:
                first_col = "Dirección" if "Dirección" in cols else ("Fuente" if "Fuente" in cols else (sorted(cols)[0] if cols else "Dirección"))
                assignment = (detected, ["URL"], [first_col], str(fpath))
        
        # 4. Try filename hint for any unmatched (direct match, no signature needed)
        if not assignment:
            hint = _hint_from_filename(fname)
            if hint:
                first_col = "Dirección" if "Dirección" in cols else ("Fuente" if "Fuente" in cols else (sorted(cols)[0] if cols else "Dirección"))
                assignment = (hint, ["URL"], [first_col], str(fpath))
        
        if assignment:
            sheet_name, headers, col_map, path = assignment
            if sheet_name not in assignments:
                assignments[sheet_name] = (headers, col_map, path)
                processed_files.add(fname)
        else:
            unmatched.append({
                "file": fname,
                "columns": sorted(cols),
                "path": str(fpath),
            })
    
    return assignments, unmatched


# ======================================================================
# Sheet ordering & Resumen
# ======================================================================

RESUMEN_SPEC = [
    ("¿Hay URLs con respuesta 4XX?", "Errores 404", "Rastreo"),
    ("¿Hay redirecciones 3XX por auditar?", "Redirecciones 3xx", "Rastreo"),
    ("¿Hay recursos HTTP (no HTTPS)?", "URLs HTTP (no HTTPS)", "Rastreo"),
    ("¿Hay bucles o cadenas de redirección?", "Bucles y cadenas redirección", "Rastreo"),
    ("URLs con caracteres no ASCII", "URLs no ASCII", "Rastreo"),
    ("Oportunidad minificar CSS", "PS Minificar CSS", "Rastreo"),
    ("Oportunidad minificar JS", "PS Minificar JS", "Rastreo"),
    ("Páginas sobredimensionadas (imágenes)", "Imágenes +100kb", "Rastreo"),
    ("Oportunidad WebP / nueva entrega imágenes", "PS Mejorar entrega imágenes", "Rastreo"),
    ("URLs sin canónica", "Canónica — Falta", "Indexación"),
    ("URLs con múltiples canónicas", "Canónica — Múltiple", "Indexación"),
    ("URLs con errores de canónica", "Canónica — Errores", "Indexación"),
    ("URLs con directiva Noindex", "Directivas — Noindex", "Indexación"),
    ("URLs con directiva Nofollow", "Directivas — Nofollow", "Indexación"),
    ("Metatítulo ausente", "Metatítulo — Falta", "On page"),
    ("Metatítulo duplicado", "Metatítulo — Duplicado", "On page"),
    ("Metatítulo múltiple", "Metatítulo — Múltiple", "On page"),
    ("Metatítulo muy largo (+60 car / +561 px)", "Metatítulo — Largo", "On page"),
    ("Metatítulo muy corto (−30 car / −200 px)", "Metatítulo — Corto", "On page"),
    ("Metadescripción ausente", "Metadescripción — Falta", "On page"),
    ("Metadescripción duplicada", "Metadescripción — Duplicada", "On page"),
    ("Metadescripción larga (+155 car / +985 px)", "Metadescripción — Larga", "On page"),
    ("Metadescripción corta (−70 car / −400 px)", "Metadescripción — Corta", "On page"),
    ("Metadescripción múltiple", "Metadescripción — Múltiple", "On page"),
    ("URLs sin H1", "H1 — Falta", "Contenido"),
    ("URLs con múltiples H1", "H1 — Múltiple", "Contenido"),
    ("URLs con H1 duplicado", "H1 — Duplicado", "Contenido"),
    ("URLs con H1 muy largo (+70 car)", "H1 — Largo (+70)", "Contenido"),
    ("URLs con poco contenido", "Poco contenido", "Contenido"),
    ("Imágenes +100kb", "Imágenes +100kb", "Contenido"),
    ("Imágenes sin ALT text", "Imágenes sin ALT text", "Contenido"),
    ("Imágenes sin atributo de tamaño", "Imágenes sin atributo tamaño", "Contenido"),
    ("Imágenes con dimensiones incorrectas", "Imágenes dim incorrectas", "Contenido"),
    ("PageSpeed − visualización de fuentes", "PS Visualización fuentes", "Indexación"),
    ("PageSpeed — LCP", "PS Solicitudes LCP", "Indexación"),
    ("PageSpeed — CLS", "PageSpeed — CLS", "Indexación"),
    ("URLs con metadatos propuestos", "Metadatos propuestos", "On page"),
]

SHEET_ORDER = [
    "Errores 404", "Detalle Errores 404",
    "Redirecciones 3xx", "Bucles y cadenas redirección",
    "URLs HTTP (no HTTPS)", "URLs no ASCII",
    "PS Minificar CSS", "PS Minificar JS",
    "PS Visualización fuentes", "PS Mejorar entrega imágenes",
    "PS Solicitudes LCP", "PageSpeed — CLS",
    "Canónica — Falta", "Canónica — Múltiple", "Canónica — Errores",
    "Directivas — Noindex", "Directivas — Nofollow",
    "Metatítulo — Falta", "Metatítulo — Duplicado", "Metatítulo — Múltiple",
    "Metatítulo — Largo", "Metatítulo — Corto", "Metatítulo = H1",
    "Metadescripción — Falta", "Metadescripción — Duplicada",
    "Metadescripción — Larga", "Metadescripción — Corta", "Metadescripción — Múltiple",
    "H1 — Falta", "H1 — Múltiple", "H1 — Duplicado", "H1 — Largo (+70)",
    "Poco contenido",
    "Imágenes +100kb", "Detalle Imágenes +100kb",
    "Imágenes sin ALT text", "Detalle Imágenes sin ALT",
    "Imágenes sin atributo tamaño", "Detalle Imágenes sin size attr",
    "Imágenes dim incorrectas", "Detalle Img dim incorrectas",
    "Metadatos propuestos",
]

PPT_HALLAZGOS = [
    (10, "Marcador de texto 2", "Errores 404"),
    (10, "Marcador de texto 4", "Redirecciones 3xx"),
    (13, "Marcador de texto 4", "Canónica — Errores"),
    (18, "Marcador de texto 4", "PS Minificar CSS"),
    (18, "Marcador de texto 6", "PS Minificar JS"),
    (19, "Marcador de texto 2", "Imágenes +100kb"),
    (22, "Marcador de texto 2", "Imágenes +100kb"),
    (22, "Marcador de texto 4", "Imágenes sin ALT text"),
    (23, "Marcador de texto 4", "Metadescripción — Duplicada"),
    (23, "Marcador de texto 6", "Poco contenido"),
    (26, "Marcador de texto 2", "H1 — Falta"),
    (26, "Marcador de texto 4", "H1 — Duplicado"),
    (26, "Marcador de texto 6", "H1 — Largo (+70)"),
    (27, "Marcador de texto 4", "Metatítulo — Corto"),
    (27, "Marcador de texto 6", "Metatítulo — Largo"),
    (28, "Marcador de texto 4", "Metadescripción — Corta"),
    (28, "Marcador de texto 6", "Metadescripción — Larga"),
]


def _reorder_sheets(wb):
    fixed = ["Matriz de entendimiento", "Resumen"]
    desired = fixed + [s for s in SHEET_ORDER if s in wb.sheetnames and s not in fixed]
    current = wb.sheetnames
    if current == desired:
        return
    for i, name in enumerate(desired):
        if i >= len(current):
            break
        if current[i] != name and name in current:
            idx = wb.sheetnames.index(name)
            wb.move_sheet(name, offset=i - idx)


# ======================================================================
# Main entry point
# ======================================================================

def build_audit(data_dir: str, xlsx_path: str, pptx_path: str = None,
                client: str = "", domain: str = "",
                skip_ppt: bool = False, skip_metadatos: bool = True,
                mapping: dict = None) -> dict:
    """
    Main function: reads data from data_dir, builds Excel + PPT.
    Auto-detects file→sheet assignments by column signatures.
    
    Returns: {"counts": {sheet: n}, "unmatched": [...], "total": N}
    """
    data_dir = Path(data_dir)
    wb = _load_wb(xlsx_path)
    counts = {}
    unmatched = []

    # 1. Detect assignments
    assignments, unmatched = _detect_assignments(str(data_dir), mapping)
    
    # 1.5 Derive sub-sheets from "Todo" element exports
    # Scan for master files (Títulos Todo, Meta Todo, H1 Todo)
    processed_elem_types = set()
    for fpath in sorted(data_dir.glob("*")):
        if fpath.suffix.lower() not in (".csv", ".ndjson", ".json"):
            continue
        if fpath.name.startswith("."):
            continue
        cols = _get_columns(fpath)
        elem_type = _detect_element_type(cols)
        if not elem_type or elem_type in processed_elem_types:
            continue
        processed_elem_types.add(elem_type)
        
        # Read ALL rows from this master file + any sibling files with same columns
        all_rows = []
        # Find all files with same element fingerprint
        fingerprint = ELEMENT_FINGERPRINTS[elem_type]
        for f2 in sorted(data_dir.glob("*")):
            if f2.suffix.lower() not in (".csv", ".ndjson", ".json"):
                continue
            f2_cols = _get_columns(f2)
            if fingerprint.issubset(f2_cols):
                rows = _read_file(f2)
                if rows:
                    all_rows.extend(rows)
        
        if not all_rows:
            continue
        
        # For images: only derive if data is mixed (Todo, not a filter export)
        if elem_type == "images":
            detected = _sample_detect(fpath)
            if not detected or detected != "__derive_images":
                continue  # Skip — this is a filtered export, not Todo
        
        # Filter by domain and dedup
        all_rows = _filter_by_domain(all_rows, domain)
        all_rows = _dedup_by_key(all_rows, "Dirección")
        
        # Derive sub-sheets
        derived = _derive_subsheets(all_rows, elem_type)
        
        for sheet_name, (headers, rows) in derived.items():
            _fill_sheet(wb, sheet_name, headers, rows)
            counts[sheet_name] = len(rows)
            # Remove from individual assignments if present (Todo overrides individual)
            assignments.pop(sheet_name, None)
        
        # Continue processing next element type (don't break)
        continue
    
    # Ensure all known sheets have a placeholder
    all_known_sheets = set()
    for sig, entries in SIGNATURES.items():
        for e in entries:
            if e["sheet"]:
                all_known_sheets.add(e["sheet"])
    all_known_sheets.update({
        "Bucles y cadenas redirección", "Metatítulo = H1", "Detalle Img dim incorrectas",
        "URLs HTTP (no HTTPS)", "Metadescripción — Falta", "H1 — Falta",
        "Canónica — Falta", "Directivas — Noindex", "Directivas — Nofollow",
        "PS Minificar CSS", "Metadescripción — Corta", "Metadescripción — Múltiple",
        "H1 — Múltiple", "Metatítulo — Múltiple", "Canónica — Errores",
        "Metadescripción — Duplicada",
    })
    # 2. Populate sheets from detected assignments
    for sheet_name, (headers, col_map, file_path) in assignments.items():
        rows_raw = _read_file(Path(file_path))
        if not rows_raw:
            _fill_sheet(wb, sheet_name, headers, [])
            counts[sheet_name] = 0
            continue

        # Filter by domain (except for certain sheets)
        if sheet_name not in ("URLs HTTP (no HTTPS)",):
            url_key = "Dirección"
            if col_map and col_map[0] and col_map[0] != "Dirección":
                url_key = col_map[0]
            rows_raw = _filter_by_domain(rows_raw, domain, key=url_key)

        # Map columns
        rows_out = []
        for r in rows_raw:
            row = []
            for key in col_map:
                if key is None:
                    row.append("")
                else:
                    val = r.get(key, "")
                    row.append(str(val) if val is not None else "")
            rows_out.append(row)
        
        _fill_sheet(wb, sheet_name, headers, rows_out)
        counts[sheet_name] = len(rows_out)

    # 3. Ensure remaining known sheets exist with "Sin hallazgos"
    for sheet_name in sorted(all_known_sheets):
        if sheet_name not in counts:
            if sheet_name == "Metadatos propuestos":
                _fill_sheet(wb, sheet_name, ["URL", "Factor a corregir", "Texto actual", "Propuesta"], [])
            elif sheet_name == "Directivas — Noindex":
                _fill_sheet(wb, sheet_name, ["URL", "Indexabilidad", "Estado de indexabilidad"], [])
            elif sheet_name == "Directivas — Nofollow":
                _fill_sheet(wb, sheet_name, ["URL", "Directivas"], [])
            else:
                _fill_sheet(wb, sheet_name, ["URL"], [])
            counts[sheet_name] = 0

    # 4. Metadatos propuestos
    if not skip_metadatos:
        metadatos_json = data_dir / "rows_metadatos.json"
        if metadatos_json.exists():
            with open(metadatos_json, encoding="utf-8") as f:
                met_rows = json.load(f)
            _fill_metadatos_propuestos(wb, met_rows)
            counts["Metadatos propuestos"] = len(met_rows)

    # 5. Rebuild Resumen
    ws_resumen = wb["Resumen"] if "Resumen" in wb.sheetnames else wb.create_sheet("Resumen")
    ws_resumen.delete_rows(1, ws_resumen.max_row)
    ws_resumen.append(["RICOPA", "FACTORES", "DESCRIPCIÓN DE FACTORES", "CANTIDAD", "VER"])
    ws_resumen.column_dimensions["A"].width = 10
    ws_resumen.column_dimensions["B"].width = 40
    ws_resumen.column_dimensions["C"].width = 50
    ws_resumen.column_dimensions["D"].width = 12
    ws_resumen.column_dimensions["E"].width = 8

    current_cat = None
    for factor_text, sheet_name, cat in RESUMEN_SPEC:
        count = counts.get(sheet_name, 0)
        row_idx = ws_resumen.max_row + 1
        ws_resumen.cell(row=row_idx, column=1, value=cat if cat != current_cat else "")
        ws_resumen.cell(row=row_idx, column=2, value=factor_text)
        ws_resumen.cell(row=row_idx, column=3, value=factor_text)
        ws_resumen.cell(row=row_idx, column=4, value=count)
        ver_cell = ws_resumen.cell(row=row_idx, column=5, value="VER")
        ver_cell.hyperlink = f"#'{sheet_name}'!A1"
        ver_cell.font = Font(color="0563C1", underline="single")
        current_cat = cat

    # 6. Reorder & save
    _reorder_sheets(wb)
    _save_wb(wb, xlsx_path)

    # 7. PowerPoint
    if not skip_ppt and pptx_path:
        excel_filename = os.path.basename(xlsx_path)
        prs = PptxPresentation(pptx_path)
        today_str = date.today().isoformat()

        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Tivit" in (run.text or ""):
                            run.text = run.text.replace("Tivit", client)

        for shape in prs.slides[1].shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if re.search(r"\d{4}-\d{2}-\d{2}", run.text or ""):
                            run.text = re.sub(r"\d{4}-\d{2}-\d{2}", today_str, run.text or "")

        for slide_idx, shape_substr, sheet_name in PPT_HALLAZGOS:
            count = counts.get(sheet_name, 0)
            slide = prs.slides[slide_idx - 1]
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape_substr.lower() not in shape.name.lower():
                    continue
                tf = shape.text_frame
                edited_here = False
                for para in tf.paragraphs:
                    for run in list(para.runs):
                        t = run.text or ""
                        if re.search(r"Hay \d+", t):
                            run.text = re.sub(r"Hay \d+", f"Hay {count}", t)
                            run.text = re.sub(r"\s*VER\s*$", "", run.text)
                            edited_here = True
                if not edited_here:
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = f"Hay {count} hallazgos."
                last_para = tf.paragraphs[-1]
                ver_run = last_para.add_run()
                ver_run.text = "  VER"
                ver_run.font.name = BRAND_FONT
                ver_run.font.size = Pt(14)
                try:
                    ver_run.font.color.rgb = RGBColor(0x05, 0x63, 0xC1)
                except Exception:
                    pass
                try:
                    ver_run.hyperlink.address = f"{excel_filename}#'{sheet_name}'!A1"
                except Exception:
                    pass

        prs.save(pptx_path)

    total = sum(counts.values())
    return {"counts": counts, "unmatched": unmatched, "total": total}
